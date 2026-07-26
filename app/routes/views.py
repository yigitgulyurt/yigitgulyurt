from flask import Blueprint, render_template, Response, url_for, current_app, request, jsonify, redirect, abort, flash, send_file, session
from flask_login import login_user, logout_user, login_required, current_user
from werkzeug.utils import secure_filename
from app.models import Project, BlogPost, StreamConfig, QrRedirect, Admin, ContactMessage, IpLog, FileShare
from app import db, limiter
from datetime import datetime, timedelta
from sqlalchemy.exc import OperationalError, DatabaseError, IntegrityError
import random
import string
import re
import os
import io
import markdown as md
import requests
import bleach
from threading import Thread
from urllib.parse import urlparse, parse_qs
from functools import lru_cache
from PIL import Image, ImageDraw, ImageFont

# --- Blueprints ---
main_bp     = Blueprint('main', __name__)
blog_bp     = Blueprint('blog', __name__)
projects_bp = Blueprint('projects', __name__)
contact_bp  = Blueprint('contact', __name__)
admin_bp    = Blueprint('admin', __name__)
og_bp       = Blueprint('og', __name__)
tools_bp    = Blueprint('tools', __name__)

# --- Helpers ---
def send_async_notification(app, name, email, subject, message):
    with app.app_context():
        send_admin_notification(name, email, subject, message)

def send_admin_notification(name, email, subject, message):
    telegram_token = current_app.config.get('TELEGRAM_TOKEN')
    admin_id       = current_app.config.get('ADMIN_TELEGRAM_ID')
    if telegram_token and admin_id:
        try:
            text = (f"📩 *Yeni İletişim Mesajı (yigitgulyurt.net.tr)*\n\n"
                    f"👤 *Gönderen:* {name}\n"
                    f"📧 *E-posta:* {email}\n"
                    f"📌 *Konu:* {subject.capitalize() if subject else 'Yok'}\n\n"
                    f"📝 *Mesaj:*\n{message}")
            requests.post(
                f"https://api.telegram.org/bot{telegram_token}/sendMessage",
                json={"chat_id": admin_id, "text": text, "parse_mode": "Markdown"},
                timeout=10,
            )
        except Exception as e:
            current_app.logger.error(f"Telegram notification error: {e}")

def generate_id(length=7):
    return ''.join(random.choices(string.ascii_letters + string.digits, k=length))

def clean_slug(text):
    text = text.lower()
    text = re.sub(r'[^a-z0-9\-_]', '', text)
    return text[:30]

def slugify(text):
    text = text.lower().strip()
    replacements = {"ç": "c", "ğ": "g", "ı": "i", "ö": "o", "ş": "s", "ü": "u"}
    for tr, en in replacements.items():
        text = text.replace(tr, en)
    text = re.sub(r'[\s]+', '-', text)
    text = re.sub(r'[^\w\-]', '', text)
    return text

def extract_slug(url):
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        if domain.startswith('www.'):
            domain = domain[4:]
        
        domain_map = {
            'youtube.com': 'yt', 'youtu.be': 'yt', 'twitch.tv': 'tv', 'vimeo.com': 'vm', 'netflix.com': 'nf',
            'twitter.com': 'tw', 'x.com': 'tw', 'instagram.com': 'ig', 'facebook.com': 'fb', 'fb.com': 'fb',
            'tiktok.com': 'tk', 'reddit.com': 'rd', 'linkedin.com': 'ln', 'pinterest.com': 'pin',
            'spotify.com': 'sp', 'open.spotify.com': 'sp', 'soundcloud.com': 'sc', 'music.apple.com': 'am',
            'github.com': 'gh', 'gitlab.com': 'gl', 'stackoverflow.com': 'so', 'codepen.io': 'cp',
            'medium.com': 'md', 'dev.to': 'dev', 'behance.net': 'be', 'dribbble.com': 'dr',
            'amazon.com': 'amz', 'amazon.com.tr': 'amz', 'wikipedia.org': 'wiki', 'discord.com': 'dc',
            'steamcommunity.com': 'st', 'steampowered.com': 'st', 't.me': 'tg', 'telegram.org': 'tg'
        }
        
        short_domain = domain_map.get(domain)
        if not short_domain:
            for d, code in domain_map.items():
                if domain.endswith('.' + d):
                    short_domain = code
                    break
        
        if not short_domain:
            short_domain = domain.split('.')[0][:10]
        
        qs = parse_qs(parsed.query)
        param_keys = ['v', 'id', 'p', 'slug', 'article', 'track', 'album', 's', 'post']
        for key in param_keys:
            if key in qs and qs[key]:
                val = qs[key][0]
                if len(val) >= 3:
                    return short_domain, clean_slug(val)
        
        path_parts = [p for p in parsed.path.split('/') if p]
        if path_parts:
            trigger_keywords = [
                'track', 'album', 'playlist', 'u', 'user', 'posts', 'p', 'watch', 
                'video', 'article', 'item', 'product', 'groups', 'community', 'channel'
            ]
            for i, part in enumerate(path_parts[:-1]):
                if part.lower() in trigger_keywords:
                    return short_domain, clean_slug(path_parts[i+1])
            
            last_segment = path_parts[-1]
            if '.' in last_segment:
                last_segment = last_segment.rsplit('.', 1)[0]
            
            if len(last_segment) >= 3:
                if last_segment.lower() in [short_domain.lower(), domain.split('.')[0].lower()]:
                    return short_domain, None
                return short_domain, clean_slug(last_segment)
                
        return short_domain, None
    except Exception:
        pass
    return "link", None

# --- Main Routes ---
@main_bp.route('/subdomain')
def subdomain():
    subdomain_list = [
        {
            "name": "Ana Site",
            "url": "https://yigitgulyurt.net.tr",
            "description": "Kişisel portfolyo ve ana site"
        },
        {
            "name": "Görsel Galerisi",
            "url": "https://yigitgulyurt.net.tr/image",
            "description": "Görsel galerisi ve dosya paylaşımı"
        },
        {
            "name": "Font Kütüphanesi",
            "url": "https://yigitgulyurt.net.tr/font",
            "description": "Özel font servisi ve font kütüphanesi"
        },
        {
            "name": "Canlı Yayın",
            "url": "https://yigitgulyurt.net.tr/canli",
            "description": "Canlı yayın platformu"
        },
        {
            "name": "CSS",
            "url": "https://css.yigitgulyurt.net.tr",
            "description": "CSS dosyaları CDN"
        },
        {
            "name": "JS",
            "url": "https://js.yigitgulyurt.net.tr",
            "description": "JavaScript dosyaları CDN"
        }
    ]
    return jsonify({
        "status": "success",
        "count": len(subdomain_list),
        "subdomain": subdomain_list
    })

@main_bp.route('/')
def index():
    featured_projects = Project.query.filter_by(featured=True).order_by(Project.order).limit(3).all()
    recent_posts = BlogPost.query.filter_by(published=True).order_by(BlogPost.created_at.desc()).limit(3).all()
    stats = {
        'projects': Project.query.count(),
        'posts': BlogPost.query.filter_by(published=True).count(),
    }
    stream_config = StreamConfig.get()
    show_stream = stream_config.show_section
    stream_live = False
    if show_stream:
        try:
            import requests as req_lib
            key = stream_config.stream_key or current_app.config.get('STREAM_KEY', '')
            r = req_lib.get('http://localhost:9997/v3/paths/list', timeout=1)
            items = r.json().get('items', [])
            path = next((p for p in items if p.get('name') == f'canli/{key}'), None)
            stream_live = bool(path and path.get('ready'))
        except Exception:
            stream_live = False

    return render_template('main/index.html',
                           featured_projects=featured_projects,
                           recent_posts=recent_posts,
                           stats=stats,
                           show_stream=show_stream,
                           stream_live=stream_live,
                           stream_config=stream_config)

@main_bp.route('/Mustafa-Kemal-Ataturk')
def ataturk():
    return render_template('main/ataturk.html')

@main_bp.route('/hakkimda')
def about():
    projects = Project.query.order_by(Project.order, Project.created_at.desc()).all()
    return render_template('main/about.html', projects=projects)

@main_bp.route('/cv')
def cv():
    projects = Project.query.order_by(Project.order, Project.created_at.desc()).all()
    return render_template('main/cv.html', projects=projects)

@main_bp.route('/lisans')
def lisans():
    license_path = os.path.join(os.path.dirname(current_app.root_path), 'LICENSE')
    try:
        with open(license_path, 'r', encoding='utf-8') as f:
            license_content = f.read()
    except Exception as e:
        current_app.logger.error(f"Lisans dosyası okunamadı: {e}")
        license_content = "Lisans dosyası bulunamadı."
    return render_template('main/license.html', license_content=license_content)

@main_bp.route('/sitemap.xml')
def sitemap():
    pages = []
    base = 'https://yigitgulyurt.net.tr'
    static_pages = [
        ('main.ataturk',   '0.8',  'monthly'),
        ('main.index',    '1.0',  'weekly'),
        ('main.about',    '0.8',  'monthly'),
        ('main.cv',       '0.7',  'monthly'),
        ('main.lisans',   '0.7',  'monthly'),
        ('projects.index','0.9',  'weekly'),
        ('blog.index',    '0.9',  'weekly'),
        ('contact.index', '0.5',  'monthly'),
        ('main.qr_okuyucu', '0.8',  'monthly'),
        ('main.file_converter', '0.8', 'monthly'),
    ]
    for endpoint, priority, changefreq in static_pages:
        pages.append({
            'loc': base + url_for(endpoint),
            'priority': priority,
            'changefreq': changefreq,
            'lastmod': datetime.utcnow().strftime('%Y-%m-%d'),
        })
    pages.append({
        'loc': 'https://yigitgulyurt.net.tr/font/',
        'priority': '0.9',
        'changefreq': 'weekly',
        'lastmod': datetime.utcnow().strftime('%Y-%m-%d'),
    })
    for p in Project.query.all():
        pages.append({
            'loc': base + url_for('projects.detail', slug=p.slug),
            'priority': '0.8',
            'changefreq': 'monthly',
            'lastmod': p.created_at.strftime('%Y-%m-%d'),
        })
    for post in BlogPost.query.filter_by(published=True).all():
        pages.append({
            'loc': base + url_for('blog.detail', slug=post.slug),
            'priority': '0.7',
            'changefreq': 'monthly',
            'lastmod': post.updated_at.strftime('%Y-%m-%d'),
        })
    xml = render_template('main/sitemap.xml', pages=pages)
    return Response(xml, mimetype='application/xml')

@main_bp.route('/api/shorten', methods=['POST'])
def shorten():
    data = request.get_json(silent=True, force=True)
    if not data: return jsonify({'error': 'JSON parse edilemedi'}), 400
    url = data.get('url', '').strip()
    if not url: return jsonify({'error': 'URL gerekli'}), 400
    existing = QrRedirect.query.filter_by(url=url).first()
    if existing:
        return jsonify({'short_id': existing.id, 'short_domain': existing.short_domain or 'r'})
    short_domain_val, slug = extract_slug(url)
    if slug:
        short_id = slug
        if QrRedirect.query.get(short_id):
            short_id = f"{slug}-{generate_id(3)}"
            while QrRedirect.query.get(short_id):
                short_id = f"{slug}-{generate_id(3)}"
    else:
        short_id = generate_id()
        while QrRedirect.query.get(short_id):
            short_id = generate_id()
    redirect_obj = QrRedirect(id=short_id, short_domain=short_domain_val, url=url)
    db.session.add(redirect_obj)
    db.session.commit()
    return jsonify({'short_id': short_id, 'short_domain': short_domain_val})

@main_bp.route('/r/<short_id>')
@main_bp.route('/r/<short_domain>/<short_id>')
def redirect_url(short_id, short_domain=None):
    obj = QrRedirect.query.get_or_404(short_id)
    obj.hit_count += 1
    db.session.commit()
    return redirect(obj.url)

@main_bp.route('/qr-okuyucu')
def qr_okuyucu():
    return render_template('qr-reader/qr-reader.html')

@main_bp.route('/robots.txt')
def robots():
    content = "User-agent: *\nAllow: /\nDisallow: /admin/\n\nSitemap: https://yigitgulyurt.net.tr/sitemap.xml\n"
    return Response(content, mimetype='text/plain')

@main_bp.route('/font-test')
def font_test():
    return render_template('main/font-test.html')

@main_bp.route('/sw.js')
def serve_sw():
    version = "1.0.0"
    sw_path = os.path.join(current_app.root_path, '..', 'sw.js')
    try:
        with open(sw_path, 'r', encoding='utf-8') as f:
            content = f.read()
        content = content.replace('${VERSION}', version)
        resp = Response(content, mimetype='application/javascript; charset=utf-8')
        resp.headers['Cache-Control'] = 'no-cache'
        return resp
    except Exception:
        return "Service Worker Not Found", 404

@main_bp.route('/offline')
def offline():
    return render_template('main/offline.html')

# --- Blog Routes ---
@blog_bp.route('/')
def index():
    posts = BlogPost.query.filter_by(published=True).order_by(BlogPost.created_at.desc()).all()
    return render_template('blog/index.html', posts=posts)

@blog_bp.route('/<slug>')
def detail(slug):
    post = BlogPost.query.filter_by(slug=slug, published=True).first_or_404()
    post.content_html = md.markdown(post.content or '', extensions=['fenced_code', 'tables'])
    return render_template('blog/detail.html', post=post)

# --- Project Routes ---
@projects_bp.route('/')
def index():
    projects = Project.query.order_by(Project.order, Project.created_at.desc()).all()
    return render_template('projects/index.html', projects=projects)

@projects_bp.route('/<slug>')
def detail(slug):
    project = Project.query.filter_by(slug=slug).first_or_404()
    
    # CSS dosyasını şablona gönder (domain'den)
    custom_css = None
    if project.custom_css_file:
        custom_css = f"https://css.yigitgulyurt.net.tr/{project.custom_css_file}"
    
    if project.content_type == 'html':
        if project.custom_html_file:
            return render_template(project.custom_html_file, project=project, custom_css=custom_css)
        else:
            project.content_html = project.content or ''
    elif project.content_type == 'template' and project.template_name:
        return render_template(project.template_name, project=project, custom_css=custom_css)
    elif project.content_type == 'markdown':
        if project.custom_markdown_file:
            # Markdown dosyasını oku ve parse et
            markdown_path = os.path.join(current_app.root_path, 'templates', project.custom_markdown_file)
            try:
                with open(markdown_path, 'r', encoding='utf-8') as f:
                    markdown_content = f.read()
                    project.content_html = md.markdown(markdown_content, extensions=['fenced_code', 'tables'])
            except FileNotFoundError:
                project.content_html = md.markdown(project.content or '', extensions=['fenced_code', 'tables'])
        else:
            project.content_html = md.markdown(project.content or '', extensions=['fenced_code', 'tables'])
    else:
        project.content_html = md.markdown(project.content or '', extensions=['fenced_code', 'tables'])
    
    return render_template('projects/detail.html', project=project, custom_css=custom_css)

# --- Contact Routes ---
@contact_bp.route('/', methods=['GET', 'POST'])
@limiter.limit("5 per hour", methods=['POST'])
def index():
    if request.method == 'POST':
        # Honeypot check
        if request.form.get('website'):
            return redirect(url_for('contact.index'))

        name    = bleach.clean(request.form.get('name', '').strip())
        email   = bleach.clean(request.form.get('email', '').strip())
        subject = bleach.clean(request.form.get('subject', '').strip())
        message = bleach.clean(request.form.get('message', '').strip())
        source  = bleach.clean(request.form.get('source', '').strip())

        if not name or not email or not message:
            flash('Lütfen zorunlu alanları doldurun.', 'error')
            return render_template('contact/index.html')

        if len(message) < 10:
            flash('Mesajınız çok kısa, lütfen biraz daha detay verin.', 'error')
            return render_template('contact/index.html')

        if len(message) > 3000:
            flash('Mesajınız çok uzun, lütfen daha kısa bir mesaj gönderin.', 'error')
            return render_template('contact/index.html')

        msg = ContactMessage(name=name, email=email, subject=subject, message=message)
        db.session.add(msg)
        db.session.commit()
        
        # Telegram Bildirimi Gönder
        app = current_app._get_current_object()
        Thread(target=send_async_notification, args=(app, name, email, subject, message)).start()

        flash('Mesajınız alındı, teşekkürler!', 'success')
        
        if source == 'file_converter':
            return redirect(url_for('main.file_converter'), code=302)
        else:
            return redirect(url_for('contact.index'))
    
    subject_from_query = bleach.clean(request.args.get('subject', '').strip())
    source_from_query = bleach.clean(request.args.get('source', '').strip())
    return render_template('contact/index.html', subject=subject_from_query, source=source_from_query)

# --- Admin Routes ---
@admin_bp.route('/giris', methods=['GET', 'POST'])
@limiter.limit("10 per minute", methods=['POST'])
@limiter.limit("50 per hour", methods=['POST'])
def login():
    if current_user.is_authenticated: return redirect(url_for('admin.index'))
    if request.method == 'POST':
        username = bleach.clean(request.form.get('username', '').strip())
        password = request.form.get('password', '')
        user = Admin.query.filter_by(username=username).first()
        if user and user.check_password(password):
            session.clear()
            login_user(user, remember=True, fresh=True)
            session.regenerate = True
            flash('Hoş geldiniz.', 'success')
            return redirect(url_for('admin.index'))
        flash('Hatalı kullanıcı adı veya şifre.', 'error')
    return render_template('admin/login.html')

@admin_bp.route('/cikis', methods=['POST'])
@login_required
def logout():
    logout_user()
    session.clear()
    flash('Güvenli şekilde çıkış yapıldı.', 'success')
    return redirect(url_for('main.index'))

@admin_bp.route('/')
@login_required
def index():
    stats = {
        'projects': Project.query.count(),
        'posts': BlogPost.query.count(),
        'messages': ContactMessage.query.filter_by(read=False).count(),
        'ip_logs': IpLog.query.count(),
    }
    return render_template('admin/index.html', stats=stats)

@admin_bp.route('/ip-logs/merge', methods=['POST'])
@login_required
def merge_ip_logs():
    from sqlalchemy import func
    # Her IP için en son id'yi bul (en güncel kaydı tutmak için)
    latest_ids_subquery = db.session.query(func.max(IpLog.id)).group_by(IpLog.ip).subquery()
    
    # Bu id'ler dışındaki tüm kayıtları sil
    deleted_count = IpLog.query.filter(~IpLog.id.in_(latest_ids_subquery)).delete(synchronize_session=False)
    db.session.commit()
    
    flash(f'{deleted_count} adet mükerrer IP kaydı başarıyla birleştirildi.', 'success')
    return redirect(url_for('admin.index'))

@admin_bp.route('/projeler')
@login_required
def projects():
    projects = Project.query.order_by(Project.order).all()
    return render_template('admin/projects.html', projects=projects)

@admin_bp.route('/projeler/yeni', methods=['GET', 'POST'])
@admin_bp.route('/projeler/<int:id>/duzenle', methods=['GET', 'POST'])
@login_required
def project_edit(id=None):
    project = Project.query.get_or_404(id) if id else Project()
    if request.method == 'POST':
        old_slug = project.slug
        project.title = request.form['title']
        project.slug = request.form.get('slug') or slugify(request.form['title'])
        project.description = request.form.get('description')
        project.tech_stack = request.form.get('tech_stack')
        project.content_type = request.form.get('content_type', 'markdown')
        project.template_name = request.form.get('template_name') or None
        
        # Kapak görseli yükle
        image_file = request.files.get('image')
        if image_file and image_file.filename:
            filename = secure_filename(image_file.filename)
            upload_dir = os.path.join(current_app.root_path, 'static', 'image', 'yigitgulyurt', 'projects')
            os.makedirs(upload_dir, exist_ok=True)
            image_file.save(os.path.join(upload_dir, filename))
            project.image = filename
        
        # HTML dosyası yükle
        html_file = request.files.get('html_file')
        if html_file and html_file.filename:
            html_filename = f"{project.slug}.html"
            html_upload_dir = os.path.join(current_app.root_path, 'templates', 'projects', 'html')
            os.makedirs(html_upload_dir, exist_ok=True)
            html_file.save(os.path.join(html_upload_dir, html_filename))
            project.custom_html_file = f"projects/html/{html_filename}"
            
            # Eski slug farklıysa eski dosyayı sil
            if old_slug and old_slug != project.slug:
                old_html_path = os.path.join(html_upload_dir, f"{old_slug}.html")
                if os.path.exists(old_html_path):
                    os.remove(old_html_path)
        
        # CSS dosyası yükle
        css_file = request.files.get('css_file')
        if css_file and css_file.filename:
            css_filename = f"{project.slug}.css"
            css_upload_dir = os.path.join(current_app.root_path, 'static', 'css', 'yigitgulyurt', 'projects')
            os.makedirs(css_upload_dir, exist_ok=True)
            css_file.save(os.path.join(css_upload_dir, css_filename))
            project.custom_css_file = f"yigitgulyurt/projects/{css_filename}"
            
            # Eski slug farklıysa eski dosyayı sil
            if old_slug and old_slug != project.slug:
                old_css_path = os.path.join(css_upload_dir, f"{old_slug}.css")
                if os.path.exists(old_css_path):
                    os.remove(old_css_path)
        
        # Markdown dosyası yükle
        markdown_file = request.files.get('markdown_file')
        if markdown_file and markdown_file.filename:
            markdown_filename = f"{project.slug}.md"
            markdown_upload_dir = os.path.join(current_app.root_path, 'templates', 'projects', 'markdown')
            os.makedirs(markdown_upload_dir, exist_ok=True)
            markdown_file.save(os.path.join(markdown_upload_dir, markdown_filename))
            project.custom_markdown_file = f"projects/markdown/{markdown_filename}"
            
            # Eski slug farklıysa eski dosyayı sil
            if old_slug and old_slug != project.slug:
                old_markdown_path = os.path.join(markdown_upload_dir, f"{old_slug}.md")
                if os.path.exists(old_markdown_path):
                    os.remove(old_markdown_path)
        
        project.content, project.live_url = request.form.get('content'), request.form.get('live_url')
        project.github_url = request.form.get('github_url')
        project.featured, project.order = bool(request.form.get('featured')), int(request.form.get('order') or 0)
        if not id: db.session.add(project)
        db.session.commit()
        flash('Proje kaydedildi.', 'success')
        return redirect(url_for('admin.projects'))
    return render_template('admin/project_edit.html', project=project)

@admin_bp.route('/projeler/<int:id>/sil', methods=['POST'])
@login_required
def project_delete(id):
    project = Project.query.get_or_404(id)
    db.session.delete(project)
    db.session.commit()
    flash('Proje silindi.', 'success')
    return redirect(url_for('admin.projects'))

@admin_bp.route('/blog')
@login_required
def blog():
    posts = BlogPost.query.order_by(BlogPost.created_at.desc()).all()
    return render_template('admin/blog.html', posts=posts)

@admin_bp.route('/blog/yeni', methods=['GET', 'POST'])
@admin_bp.route('/blog/<int:id>/duzenle', methods=['GET', 'POST'])
@login_required
def post_edit(id=None):
    post = BlogPost.query.get_or_404(id) if id else BlogPost()
    if request.method == 'POST':
        post.title = request.form['title']
        post.slug = request.form.get('slug') or slugify(request.form['title'])
        post.summary, post.content = request.form.get('summary'), request.form.get('content')
        post.published = bool(request.form.get('published'))
        if not id: db.session.add(post)
        db.session.commit()
        flash('Yazı kaydedildi.', 'success')
        return redirect(url_for('admin.blog'))
    return render_template('admin/post_edit.html', post=post)

@admin_bp.route('/blog/<int:id>/sil', methods=['POST'])
@login_required
def post_delete(id):
    post = BlogPost.query.get_or_404(id)
    db.session.delete(post)
    db.session.commit()
    flash('Yazı silindi.', 'success')
    return redirect(url_for('admin.blog'))

@admin_bp.route('/preview', methods=['POST'])
@login_required
def preview():
    text = request.form.get('content', '')
    html = md.markdown(text, extensions=['fenced_code', 'tables'])
    return jsonify({'html': html})

@admin_bp.route('/yayin', methods=['GET', 'POST'])
@login_required
def stream_config():
    cfg = StreamConfig.get()
    if request.method == 'POST':
        new_key = request.form.get('stream_key', '').strip()
        if new_key: cfg.stream_key = new_key
        cfg.show_section = bool(request.form.get('show_section'))
        cfg.title, cfg.subtitle = request.form.get('title', '').strip() or 'Canlı Yayın', request.form.get('subtitle', '').strip()
        db.session.commit()
        flash('Yayın ayarları güncellendi.', 'success')
        return redirect(url_for('admin.stream_config'))
    return render_template('admin/stream_config.html', cfg=cfg)

@admin_bp.route('/izleyiciler')
@login_required
def stream_viewers():
    import json, time
    viewers = []
    try:
        import redis as redis_lib
        r = redis_lib.from_url(os.environ.get('REDIS_URL', 'redis://localhost:6379/0'), decode_responses=True)
        keys = r.keys('yg_viewer:*')
        now = int(time.time())
        for key in keys:
            raw = r.get(key)
            if not raw: continue
            try:
                data = json.loads(raw)
                data['since'], data['ttl'] = now - data.get('last_seen', now), r.ttl(key)
                viewers.append(data)
            except Exception: pass
        viewers.sort(key=lambda x: x.get('last_seen', 0), reverse=True)
    except Exception: viewers = []
    return render_template('admin/stream_viewers.html', viewers=viewers)

def split_long_lines(text, max_len=80):
    lines = []
    for line in text.split('\n'):
        if len(line) <= max_len:
            lines.append(line)
        else:
            words = line.split()
            current_line = ''
            for word in words:
                while len(word) > max_len:
                    if current_line:
                        lines.append(current_line)
                        current_line = ''
                    lines.append(word[:max_len])
                    word = word[max_len:]
                if len(current_line) + len(word) + 1 <= max_len:
                    current_line = (current_line + ' ' + word) if current_line else word
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            if current_line:
                lines.append(current_line)
    return lines

@admin_bp.route('/mesajlar')
@login_required
def messages():
    msgs = ContactMessage.query.order_by(ContactMessage.created_at.desc()).all()
    return render_template('admin/messages.html', messages=msgs)

@admin_bp.route('/mesajlar/<int:id>')
@login_required
def message_detail(id):
    msg = ContactMessage.query.get_or_404(id)
    msg.read = True
    db.session.commit()
    formatted_message = split_long_lines(msg.message)
    return render_template('admin/message_detail.html', message=msg, formatted_message=formatted_message)

@admin_bp.route('/og')
@login_required
def og():
    return render_template('admin/showcase/og.html')

@admin_bp.route('/og/yigitgulyurt.net.tr')
@login_required
def yigitgulyurt_og():
    return render_template('admin/showcase/yigitgulyurt.html')

@admin_bp.route('/og/cagrivakti.com.tr')
@login_required
def cagrivakti_og():
    return render_template('admin/showcase/cagrivakti.html')

# --- Personal File Manager ---
def allowed_file(filename):
    allowed = current_app.config.get('ALLOWED_EXTENSIONS', set())
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed

def human_readable_size(num_bytes):
    for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
        if abs(num_bytes) < 1024.0:
            return f"{num_bytes:.2f} {unit}"
        num_bytes /= 1024.0
    return f"{num_bytes:.2f} PB"

def sanitize_relative_path(user_path):
    if not user_path:
        return ''
    user_path = user_path.replace('\\', '/').lstrip('/')
    user_path = re.sub(r'\.\.+', '.', user_path)
    user_path = user_path.lstrip('./')
    parts = []
    for p in user_path.split('/'):
        if not p or p in ('.', '..'):
            continue
        parts.append(secure_filename(p) if '.' not in p else p)
    return '/'.join(parts)

@admin_bp.route('/dosyalar')
@admin_bp.route('/dosyalar/<path:subpath>')
@login_required
@limiter.limit("120 per minute")
def file_manager(subpath=''):
    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    safe_sub = sanitize_relative_path(subpath)
    current_dir = os.path.join(base_dir, safe_sub) if safe_sub else base_dir
    current_dir = os.path.realpath(current_dir)
    base_real = os.path.realpath(base_dir)
    if not current_dir.startswith(base_real):
        flash('Geçersiz klasör yolu.', 'error')
        return redirect(url_for('admin.file_manager'))
    if not os.path.isdir(current_dir):
        flash('Klasör bulunamadı.', 'error')
        return redirect(url_for('admin.file_manager'))
    files = []
    dirs = []
    try:
        with os.scandir(current_dir) as it:
            for entry in it:
                try:
                    stat = entry.stat()
                    item = {
                        'name': entry.name,
                        'mtime': datetime.fromtimestamp(stat.st_mtime),
                        'size': stat.st_size,
                        'size_human': human_readable_size(stat.st_size),
                        'path': (safe_sub + '/' + entry.name).lstrip('/') if safe_sub else entry.name,
                    }
                    if entry.is_dir():
                        dirs.append(item)
                    elif entry.is_file():
                        item['ext'] = entry.name.rsplit('.', 1)[1].lower() if '.' in entry.name else ''
                        files.append(item)
                except (OSError, PermissionError):
                    continue
    except (OSError, PermissionError) as e:
        flash(f'Klasör okunamadı: {e}', 'error')
    dirs.sort(key=lambda x: x['name'].lower())
    files.sort(key=lambda x: x['name'].lower())
    breadcrumb = []
    acc = ''
    for part in ([p for p in safe_sub.split('/') if p] if safe_sub else []):
        acc = (acc + '/' + part).lstrip('/')
        breadcrumb.append({'name': part, 'path': acc})
    parent = ''
    if safe_sub:
        parent_parts = safe_sub.split('/')[:-1]
        parent = '/'.join(parent_parts)
    total_size = sum(f['size'] for f in files)
    return render_template(
        'admin/file_manager.html',
        dirs=dirs,
        files=files,
        breadcrumb=breadcrumb,
        current_path=safe_sub,
        parent_path=parent,
        total_files=len(files),
        total_dirs=len(dirs),
        total_size_human=human_readable_size(total_size),
    )

@admin_bp.route('/dosyalar/yukle', methods=['POST'])
@login_required
@limiter.limit("30 per minute")
def file_upload():
    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    target_sub = sanitize_relative_path(request.form.get('path', ''))
    target_dir = os.path.join(base_dir, target_sub) if target_sub else base_dir
    target_dir = os.path.realpath(target_dir)
    base_real = os.path.realpath(base_dir)
    if not target_dir.startswith(base_real):
        return jsonify({'status': 'error', 'message': 'Geçersiz hedef klasör.'}), 400
    os.makedirs(target_dir, exist_ok=True)
    uploaded_files = request.files.getlist('files')
    if not uploaded_files or all(not f.filename for f in uploaded_files):
        flash('Yüklenecek dosya seçilmedi.', 'error')
        return redirect(url_for('admin.file_manager', subpath=target_sub or None))
    results = []
    for f in uploaded_files:
        if not f or not f.filename:
            continue
        original = f.filename
        filename = secure_filename(original) or 'dosya'
        if not allowed_file(filename):
            results.append({'name': original, 'status': 'error', 'message': f'İzin verilmeyen uzantı: {filename.rsplit(".", 1)[1].lower() if "." in filename else "yok"}'})
            continue
        target_path = os.path.join(target_dir, filename)
        counter = 1
        name, ext = os.path.splitext(filename)
        while os.path.exists(target_path):
            filename = f"{name}_{counter}{ext}"
            target_path = os.path.join(target_dir, filename)
            counter += 1
        try:
            size = 0
            chunk_size = 8 * 1024 * 1024
            with open(target_path, 'wb') as out:
                while True:
                    chunk = f.stream.read(chunk_size)
                    if not chunk:
                        break
                    size += len(chunk)
                    out.write(chunk)
            results.append({'name': original, 'status': 'ok', 'saved_as': filename, 'size_human': human_readable_size(size)})
        except Exception as e:
            results.append({'name': original, 'status': 'error', 'message': str(e)})
    ok_count = sum(1 for r in results if r['status'] == 'ok')
    err_count = len(results) - ok_count
    if ok_count:
        flash(f'{ok_count} dosya başarıyla yüklendi.', 'success')
    if err_count:
        flash(f'{err_count} dosya yüklenemedi. Detay: ' + ', '.join(r['message'] for r in results if r['status'] == 'error'), 'error')
    return redirect(url_for('admin.file_manager', subpath=target_sub or None))

@admin_bp.route('/dosyalar/indir/<path:filepath>')
@login_required
@limiter.limit("120 per minute")
def file_download(filepath):
    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    safe = sanitize_relative_path(filepath)
    target = os.path.realpath(os.path.join(base_dir, safe))
    base_real = os.path.realpath(base_dir)
    if not target.startswith(base_real) or not os.path.isfile(target):
        abort(404)
    return send_file(target, as_attachment=True, download_name=os.path.basename(target))

@admin_bp.route('/dosyalar/sil/<path:filepath>', methods=['POST'])
@login_required
@limiter.limit("60 per minute")
def file_delete(filepath):
    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    safe = sanitize_relative_path(filepath)
    target = os.path.realpath(os.path.join(base_dir, safe))
    base_real = os.path.realpath(base_dir)
    if not target.startswith(base_real):
        flash('Geçersiz dosya yolu.', 'error')
        return redirect(url_for('admin.file_manager'))
    parent_sub = ''
    if safe:
        parts = safe.split('/')
        parent_sub = '/'.join(parts[:-1])
    try:
        if os.path.isfile(target):
            os.remove(target)
            flash(f'{os.path.basename(target)} silindi.', 'success')
        elif os.path.isdir(target):
            import shutil
            shutil.rmtree(target)
            flash(f'Klasör silindi: {os.path.basename(target)}', 'success')
        else:
            flash('Hedef bulunamadı.', 'error')
    except Exception as e:
        flash(f'Silme başarısız: {e}', 'error')
    return redirect(url_for('admin.file_manager', subpath=parent_sub or None))

@admin_bp.route('/dosyalar/klasor-olustur', methods=['POST'])
@login_required
@limiter.limit("30 per minute")
def file_mkdir():
    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    target_sub = sanitize_relative_path(request.form.get('path', ''))
    folder_name = secure_filename(request.form.get('folder_name', '').strip())
    if not folder_name:
        flash('Klasör adı gerekli.', 'error')
        return redirect(url_for('admin.file_manager', subpath=target_sub or None))
    target_dir = os.path.join(base_dir, target_sub) if target_sub else base_dir
    target_dir = os.path.realpath(target_dir)
    base_real = os.path.realpath(base_dir)
    if not target_dir.startswith(base_real):
        flash('Geçersiz hedef klasör.', 'error')
        return redirect(url_for('admin.file_manager'))
    new_dir = os.path.join(target_dir, folder_name)
    try:
        os.makedirs(new_dir, exist_ok=False)
        flash(f'"{folder_name}" klasörü oluşturuldu.', 'success')
    except FileExistsError:
        flash('Bu isimde bir klasör zaten var.', 'error')
    except Exception as e:
        flash(f'Oluşturma başarısız: {e}', 'error')
    return redirect(url_for('admin.file_manager', subpath=target_sub or None))


# --- Resumable (Chunked) Upload Endpoints ---

def _ensure_partial_dir():
    partial_dir = current_app.config['PERSONAL_UPLOAD_PARTIAL_FOLDER']
    os.makedirs(partial_dir, exist_ok=True)
    return partial_dir

def _file_identifier(filename, file_size, last_modified=None):
    """Dosya için benzersiz bir tanımlayıcı üret (chunk'ları eşleştirmek için)"""
    import hashlib
    base_str = f"{filename}|{file_size}|{last_modified or 0}"
    return hashlib.sha256(base_str.encode('utf-8')).hexdigest()[:32]

def _partial_path(file_id, chunk_index=None):
    partial_dir = _ensure_partial_dir()
    if chunk_index is None:
        return os.path.join(partial_dir, file_id)
    return os.path.join(partial_dir, f"{file_id}.{chunk_index:010d}")

@admin_bp.route('/dosyalar/yukle-durum', methods=['POST'])
@login_required
@limiter.limit("120 per minute")
def file_upload_status():
    """Mevcut yükleme durumunu döndür - frontend'in kaldığı yerden devam etmesi için"""
    data = request.get_json(silent=True) or request.form
    filename = data.get('filename', '').strip()
    file_size = int(data.get('file_size', 0) or 0)
    last_modified = data.get('last_modified')

    if not filename or file_size <= 0:
        return jsonify({'status': 'error', 'message': 'Eksik parametreler'}), 400

    file_id = _file_identifier(filename, file_size, last_modified)
    partial_dir = _ensure_partial_dir()

    uploaded_bytes = 0
    completed_chunks = []

    if os.path.isdir(partial_dir):
        for entry in os.listdir(partial_dir):
            if entry.startswith(file_id + '.'):
                try:
                    idx = int(entry.split('.')[-1])
                    chunk_path = os.path.join(partial_dir, entry)
                    sz = os.path.getsize(chunk_path)
                    uploaded_bytes += sz
                    completed_chunks.append(idx)
                except (ValueError, OSError):
                    continue

    return jsonify({
        'status': 'ok',
        'file_id': file_id,
        'uploaded_bytes': uploaded_bytes,
        'file_size': file_size,
        'completed_chunks': sorted(completed_chunks),
        'progress_pct': round((uploaded_bytes / file_size) * 100, 2) if file_size else 0,
    })

@admin_bp.route('/dosyalar/yukle-chunk', methods=['POST'])
@login_required
@limiter.limit("600 per minute")
def file_upload_chunk():
    """Tek bir chunk'ı .partials klasörüne kaydet"""
    file_id = request.form.get('file_id', '').strip()
    chunk_index = int(request.form.get('chunk_index', 0) or 0)
    chunk_total = int(request.form.get('chunk_total', 1) or 1)
    filename = request.form.get('filename', '').strip()
    file_size = int(request.form.get('file_size', 0) or 0)

    chunk_file = request.files.get('chunk')
    if not chunk_file or not file_id:
        return jsonify({'status': 'error', 'message': 'Eksik parametreler'}), 400

    if not allowed_file(filename):
        return jsonify({'status': 'error', 'message': f'İzin verilmeyen uzantı: {filename.rsplit(".", 1)[1].lower() if "." in filename else "yok"}'}), 400

    chunk_path = _partial_path(file_id, chunk_index)
    try:
        size = 0
        with open(chunk_path, 'wb') as out:
            while True:
                buf = chunk_file.stream.read(4 * 1024 * 1024)
                if not buf:
                    break
                size += len(buf)
                out.write(buf)
        return jsonify({
            'status': 'ok',
            'chunk_index': chunk_index,
            'chunk_size': size,
        })
    except Exception as e:
        try:
            if os.path.exists(chunk_path):
                os.remove(chunk_path)
        except OSError:
            pass
        return jsonify({'status': 'error', 'message': str(e)}), 500

@admin_bp.route('/dosyalar/yukle-bitir', methods=['POST'])
@login_required
@limiter.limit("60 per minute")
def file_upload_finalize():
    """Tüm chunk'ları birleştir, hedef klasöre taşı, partial'ları temizle"""
    data = request.get_json(silent=True) or request.form
    file_id = data.get('file_id', '').strip()
    filename = data.get('filename', '').strip()
    file_size = int(data.get('file_size', 0) or 0)
    chunk_total = int(data.get('chunk_total', 0) or 0)
    target_sub = sanitize_relative_path(data.get('path', ''))

    if not file_id or not filename or chunk_total <= 0:
        return jsonify({'status': 'error', 'message': 'Eksik parametreler'}), 400

    base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
    target_dir = os.path.join(base_dir, target_sub) if target_sub else base_dir
    target_dir = os.path.realpath(target_dir)
    base_real = os.path.realpath(base_dir)
    if not target_dir.startswith(base_real):
        return jsonify({'status': 'error', 'message': 'Geçersiz hedef klasör'}), 400
    os.makedirs(target_dir, exist_ok=True)

    if not allowed_file(filename):
        return jsonify({'status': 'error', 'message': f'İzin verilmeyen uzantı: {filename.rsplit(".", 1)[1].lower() if "." in filename else "yok"}'}), 400

    chunk_paths = []
    for i in range(chunk_total):
        p = _partial_path(file_id, i)
        if not os.path.isfile(p):
            return jsonify({'status': 'error', 'message': f'Eksik chunk: {i}'}), 400
        chunk_paths.append(p)

    secure_name = secure_filename(filename) or 'dosya'
    name, ext = os.path.splitext(secure_name)
    counter = 1
    final_filename = secure_name
    final_path = os.path.join(target_dir, final_filename)
    while os.path.exists(final_path):
        final_filename = f"{name}_{counter}{ext}"
        final_path = os.path.join(target_dir, final_filename)
        counter += 1

    try:
        total_written = 0
        with open(final_path, 'wb') as out:
            for p in chunk_paths:
                with open(p, 'rb') as cf:
                    while True:
                        buf = cf.read(8 * 1024 * 1024)
                        if not buf:
                            break
                        out.write(buf)
                        total_written += len(buf)
        for p in chunk_paths:
            try: os.remove(p)
            except OSError: pass
        meta_path = _partial_path(file_id)
        try:
            if os.path.exists(meta_path): os.remove(meta_path)
        except OSError:
            pass
        return jsonify({
            'status': 'ok',
            'filename': final_filename,
            'saved_path': (target_sub + '/' + final_filename).lstrip('/') if target_sub else final_filename,
            'size_bytes': total_written,
            'size_human': human_readable_size(total_written),
        })
    except Exception as e:
        try:
            if os.path.exists(final_path):
                os.remove(final_path)
        except OSError:
            pass
        return jsonify({'status': 'error', 'message': str(e)}), 500

@admin_bp.route('/dosyalar/yukle-iptal', methods=['POST'])
@login_required
@limiter.limit("60 per minute")
def file_upload_cancel():
    """Kısmi yüklemeyi iptal et - chunk'ları sil"""
    data = request.get_json(silent=True) or request.form
    file_id = data.get('file_id', '').strip()
    if not file_id:
        return jsonify({'status': 'error', 'message': 'file_id gerekli'}), 400
    partial_dir = _ensure_partial_dir()
    deleted = 0
    try:
        for entry in os.listdir(partial_dir):
            if entry.startswith(file_id):
                p = os.path.join(partial_dir, entry)
                try:
                    os.remove(p)
                    deleted += 1
                except OSError:
                    pass
    except OSError as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500
    return jsonify({'status': 'ok', 'deleted_chunks': deleted})


# --- File Share Endpoints ---

def _share_db_err(e):
    """Veritabanı hatalarını (özellikle eksik migration/tablo) kullanıcı dostu mesaja çevirir."""
    msg = str(e)
    low = msg.lower()
    if 'no such table' in low or 'relation' in low and 'does not exist' in low:
        return ('Paylaşım özelliği için veritabanı tablosu oluşturulmamış. '
                'Sunucuda `flask db upgrade` komutunu çalıştırarak migrationları uygulayın.')
    if 'column' in low and ('does not exist' in low or 'no such column' in low):
        return f'Veritabanı şeması güncel değil: {msg[:200]} · Sunucuda `flask db upgrade` çalıştırın.'
    return f'Veritabanı hatası: {msg[:220]}'


def _generate_share_token():
    length = current_app.config.get('SHARE_TOKEN_LENGTH', 28)
    import secrets
    alphabet = string.ascii_letters + string.digits
    while True:
        token = ''.join(secrets.choice(alphabet) for _ in range(length))
        try:
            if not FileShare.query.filter_by(token=token).first():
                return token
        except (OperationalError, DatabaseError):
            # DB hatasıysa çakışma kontrolünü atla, nadir durumda token çakışması olursa IntegrityError döner
            return token

@admin_bp.route('/dosyalar/paylas', methods=['POST'])
@login_required
@limiter.limit("60 per minute")
def file_share_create():
    try:
        data = request.get_json(silent=True) or request.form
        filepath = sanitize_relative_path(data.get('filepath', ''))
        if not filepath:
            return jsonify({'status': 'error', 'message': 'Dosya yolu gerekli'}), 400

        base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
        target = os.path.realpath(os.path.join(base_dir, filepath))
        base_real = os.path.realpath(base_dir)
        if not target.startswith(base_real) or not os.path.isfile(target):
            return jsonify({'status': 'error', 'message': 'Dosya bulunamadı'}), 404

        try:
            stat = os.stat(target)
        except OSError as e:
            return jsonify({'status': 'error', 'message': str(e)}), 500

        password = data.get('password') or None
        max_downloads = data.get('max_downloads')
        if max_downloads:
            try:
                max_downloads = int(max_downloads)
                if max_downloads <= 0:
                    max_downloads = None
            except (ValueError, TypeError):
                max_downloads = None
        else:
            max_downloads = None

        expires_hours = data.get('expires_hours')
        expires_at = None
        if expires_hours:
            try:
                h = float(expires_hours)
                if h > 0:
                    expires_at = datetime.utcnow() + timedelta(hours=h)
            except (ValueError, TypeError):
                expires_at = None

        share = FileShare(
            token=_generate_share_token(),
            file_path=filepath,
            file_name=os.path.basename(target),
            file_size=stat.st_size,
            max_downloads=max_downloads,
            expires_at=expires_at,
            created_by_id=current_user.id,
        )
        if password:
            share.set_password(password)

        db.session.add(share)
        db.session.commit()

        share_url = url_for('admin.file_share_view', token=share.token, _external=True)

        return jsonify({
            'status': 'ok',
            'share': {
                'id': share.id,
                'token': share.token,
                'url': share_url,
                'file_name': share.file_name,
                'file_size_human': human_readable_size(share.file_size or 0),
                'has_password': bool(share.password_hash),
                'max_downloads': share.max_downloads,
                'download_count': share.download_count or 0,
                'expires_at': share.expires_at.isoformat() if share.expires_at else None,
                'time_left': share.time_left_text(),
                'is_active': share.is_active,
                'created_at': share.created_at.isoformat() if share.created_at else None,
            }
        })
    except (OperationalError, DatabaseError, IntegrityError) as e:
        db.session.rollback()
        return jsonify({'status': 'error', 'message': _share_db_err(e)}), 503
    except Exception as e:
        db.session.rollback()
        return jsonify({'status': 'error', 'message': f'Paylaşım oluşturulamadı: {e}'}), 500

@admin_bp.route('/dosyalar/paylasimlar', methods=['GET'])
@login_required
@limiter.limit("120 per minute")
def file_share_list():
    """Tüm paylaşımları listele - ya da belirli bir filepath için"""
    try:
        filepath = request.args.get('filepath')
        query = FileShare.query.filter_by(created_by_id=current_user.id)
        if filepath:
            safe = sanitize_relative_path(filepath)
            query = query.filter_by(file_path=safe)
        shares = query.order_by(FileShare.created_at.desc()).all()
        result = []
        for s in shares:
            result.append({
                'id': s.id,
                'token': s.token,
                'url': url_for('admin.file_share_view', token=s.token, _external=True),
                'file_path': s.file_path,
                'file_name': s.file_name,
                'file_size_human': human_readable_size(s.file_size or 0),
                'has_password': bool(s.password_hash),
                'max_downloads': s.max_downloads,
                'download_count': s.download_count or 0,
                'expires_at': s.expires_at.isoformat() if s.expires_at else None,
                'time_left': s.time_left_text(),
                'is_active': s.is_active,
                'created_at': s.created_at.isoformat() if s.created_at else None,
            })
        return jsonify({'status': 'ok', 'count': len(result), 'shares': result})
    except (OperationalError, DatabaseError) as e:
        return jsonify({'status': 'error', 'message': _share_db_err(e), 'count': 0, 'shares': []}), 503
    except Exception as e:
        return jsonify({'status': 'error', 'message': f'Paylaşımlar listelenemedi: {e}', 'count': 0, 'shares': []}), 500

@admin_bp.route('/dosyalar/paylasim/<int:share_id>/sil', methods=['POST'])
@login_required
@limiter.limit("60 per minute")
def file_share_delete(share_id):
    try:
        share = FileShare.query.get_or_404(share_id)
        if share.created_by_id != current_user.id:
            abort(403)
        db.session.delete(share)
        db.session.commit()
        return jsonify({'status': 'ok', 'message': 'Paylaşım silindi'})
    except (OperationalError, DatabaseError, IntegrityError) as e:
        db.session.rollback()
        return jsonify({'status': 'error', 'message': _share_db_err(e)}), 503
    except Exception as e:
        db.session.rollback()
        return jsonify({'status': 'error', 'message': f'Silme başarısız: {e}'}), 500

# Public share endpoints
@main_bp.route('/p/<token>', methods=['GET', 'POST'])
def file_share_view(token):
    try:
        share = FileShare.query.filter_by(token=token).first()
        if not share:
            abort(404)

        error = None
        if request.method == 'POST':
            password = request.form.get('password', '')
            if not share.check_password(password):
                error = 'Hatalı şifre'
            else:
                session['share_unlock_' + token] = True
                return redirect(url_for('admin.file_share_view', token=token))

        needs_password = bool(share.password_hash) and not session.get('share_unlock_' + token)
        if not share.is_active:
            return render_template('admin/file_share_expired.html', share=share, error=error), 410

        base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
        target = os.path.realpath(os.path.join(base_dir, share.file_path))
        base_real = os.path.realpath(base_dir)
        if not target.startswith(base_real) or not os.path.isfile(target):
            return render_template('admin/file_share_expired.html', share=share, error='Dosya bulunamadı'), 404

        return render_template('admin/file_share_view.html',
                               share=share,
                               file_size_human=human_readable_size(share.file_size or 0),
                               needs_password=needs_password,
                               error=error)
    except (OperationalError, DatabaseError):
        try:
            return render_template('admin/file_share_expired.html',
                                   share=None, error='Paylaşım özelliği henüz hazır değil. Sunucu yöneticisine `flask db upgrade` komutunu çalıştırmasını söyleyin.'), 503
        except Exception:
            abort(503)

@main_bp.route('/p/<token>/indir')
def file_share_download(token):
    try:
        share = FileShare.query.filter_by(token=token).first()
        if not share:
            abort(404)
        if not share.is_active:
            abort(410)
        if share.password_hash and not session.get('share_unlock_' + token):
            return redirect(url_for('admin.file_share_view', token=token))

        base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
        target = os.path.realpath(os.path.join(base_dir, share.file_path))
        base_real = os.path.realpath(base_dir)
        if not target.startswith(base_real) or not os.path.isfile(target):
            abort(404)

        share.download_count = (share.download_count or 0) + 1
        db.session.commit()

        return send_file(target, as_attachment=True, download_name=share.file_name)
    except (OperationalError, DatabaseError, IntegrityError) as e:
        db.session.rollback()
        abort(503)
    except Exception as e:
        db.session.rollback()
        abort(500)

@main_bp.route('/p/<token>/dosya')
def file_share_raw(token):
    """Dosyayı doğrudan gönder (resim/video/markdown/kod için önizleme - indirme sayacını arttırmaz)"""
    try:
        share = FileShare.query.filter_by(token=token).first()
        if not share or not share.is_active:
            abort(404)
        if share.password_hash and not session.get('share_unlock_' + token):
            abort(403)

        base_dir = current_app.config['PERSONAL_UPLOAD_FOLDER']
        target = os.path.realpath(os.path.join(base_dir, share.file_path))
        base_real = os.path.realpath(base_dir)
        if not target.startswith(base_real) or not os.path.isfile(target):
            abort(404)

        ext = (os.path.splitext(target)[1].lstrip('.') or '').lower()
        text_exts = {'txt','md','markdown','log','csv','json','xml','yml','yaml','toml','ini','cfg','env','html','htm','css','js','ts','py','sql','sh','bash','bat','ps1','c','cpp','h','hpp','rs','go','java','kt','rb','php','pl','lua','vim','conf','nginx','dockerfile','makefile'}
        if ext in text_exts:
            try:
                mime_map = {
                    'html': 'text/plain; charset=utf-8', 'htm': 'text/plain; charset=utf-8',
                    'js': 'text/plain; charset=utf-8', 'ts': 'text/plain; charset=utf-8',
                    'css': 'text/css; charset=utf-8', 'json': 'application/json; charset=utf-8',
                    'xml': 'application/xml; charset=utf-8', 'svg': 'image/svg+xml; charset=utf-8',
                }
                mimetype = mime_map.get(ext, 'text/plain; charset=utf-8')
                response = send_file(target, as_attachment=False, download_name=share.file_name, mimetype=mimetype)
            except Exception:
                response = send_file(target, as_attachment=False, download_name=share.file_name)
        else:
            response = send_file(target, as_attachment=False, download_name=share.file_name)
        response.headers.setdefault('X-Content-Type-Options', 'nosniff')
        return response
    except (OperationalError, DatabaseError):
        abort(503)
    except Exception as e:
        abort(500)


# --- OG Routes ---
THEMES = {
    'default': {'bg': '#0d0d0d', 'accent': '#4ade80', 'accent2': '#60a5fa', 'text': '#e2e2e2', 'text2': '#888888'},
    'live':    {'bg': '#0d0d0d', 'accent': '#f87171', 'accent2': '#60a5fa', 'text': '#e2e2e2', 'text2': '#888888'},
    'ataturk': {'bg': '#080808', 'accent': '#e30a17', 'accent2': '#c5a059', 'text': '#f0f0f0', 'text2': '#999999'},
    'blog':    {'bg': '#0d0d0d', 'accent': '#fb923c', 'accent2': '#fcd34d', 'text': '#e2e2e2', 'text2': '#888888'},
    'project': {'bg': '#0d0d0d', 'accent': '#818cf8', 'accent2': '#22d3ee', 'text': '#e2e2e2', 'text2': '#888888'},
    'contact': {'bg': '#051010', 'accent': '#10b981', 'accent2': '#2dd4bf', 'text': '#f9fafb', 'text2': '#94a3b8'},
    'about':   {'bg': '#0f172a', 'accent': '#a78bfa', 'accent2': '#f472b6', 'text': '#f8fafc', 'text2': '#94a3b8'},
    'cv':      {'bg': '#0f172a', 'accent': '#38bdf8', 'accent2': '#94a3b8', 'text': '#f1faee', 'text2': '#a8dadc'},
    'main':    {'bg': '#111827', 'accent': '#f59e0b', 'accent2': '#fbbf24', 'text': '#f9fafb', 'text2': '#9ca3af'},
    'qr':      {'bg': '#000000', 'accent': '#ffffff', 'accent2': '#888888', 'text': '#ffffff', 'text2': '#cccccc'},
}

FONT_BOLD = 'app/static/fonts/JetBrainsMonoNerdFont/JetBrainsMonoNerdFont-Bold.ttf'
FONT_REG  = 'app/static/fonts/JetBrainsMonoNerdFont/JetBrainsMonoNerdFont-Regular.ttf'
W, H, PAD, MARGIN = 1200, 630, 64, 40
DOM_PAD_X, DOM_PAD_Y, DOM_FONT_SZ, PROMPT_SZ, SUBTITLE_SZ = 16, 10, 20, 22, 26
BRACKET_LW, TITLE_SIZES = 2, (72, 58, 46, 36, 28)

def _hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def _load_font(path, size):
    try: return ImageFont.truetype(path, size)
    except Exception: return ImageFont.load_default()

def _fit_title_font(draw, text, max_width):
    for size in TITLE_SIZES:
        font = _load_font(FONT_BOLD, size)
        if draw.textbbox((0, 0), text, font=font)[2] < max_width: return font, size
    return _load_font(FONT_BOLD, TITLE_SIZES[-1]), TITLE_SIZES[-1]

def _draw_bracket(draw, x1, y1, x2, y2, color):
    arm, lw = max(20, min(36, int((x2 - x1) * 0.18))), BRACKET_LW
    draw.rectangle([x1, y1, x1 + arm, y1 + lw], fill=color)
    draw.rectangle([x1, y1, x1 + lw, y1 + arm], fill=color)
    draw.rectangle([x2 - arm, y1, x2, y1 + lw], fill=color)
    draw.rectangle([x2 - lw, y1, x2, y1 + arm], fill=color)
    draw.rectangle([x1, y2 - lw, x1 + arm, y2], fill=color)
    draw.rectangle([x1, y2 - arm, x1 + lw, y2], fill=color)
    draw.rectangle([x2 - arm, y2 - lw, x2, y2], fill=color)
    draw.rectangle([x2 - lw, y2 - arm, x2, y2], fill=color)

def _draw_subtitle_multiline(draw, text, x, y, font, color, max_width, line_spacing=12):
    if '|' in text:
        line_h = draw.textbbox((0, 0), 'A', font=font)[3] + line_spacing
        for i, line in enumerate(text.split('|')): draw.text((x, y + i * line_h), line.strip(), font=font, fill=color)
        return
    words, line, cy = text.split(' '), '', y
    for word in words:
        test = (line + ' ' + word).strip()
        if draw.textbbox((0, 0), test, font=font)[2] <= max_width: line = test
        else:
            if line:
                draw.text((x, cy), line, font=font, fill=color)
                cy += draw.textbbox((0, 0), line, font=font)[3] + line_spacing
            line = word
    if line: draw.text((x, cy), line, font=font, fill=color)

def make_og(title, subtitle, theme, prompt, domain):
    t = THEMES.get(theme, THEMES['default'])
    img = Image.new('RGB', (W, H), _hex_to_rgb(t['bg']))
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, W, 3], fill=_hex_to_rgb(t['accent']))
    f_prompt = _load_font(FONT_REG, PROMPT_SZ)
    d.text((PAD, 48), prompt, font=f_prompt, fill=_hex_to_rgb(t['accent']))
    max_title_w = W - PAD * 2
    f_title, title_size = _fit_title_font(d, title, max_title_w)
    title_y = H // 2 - title_size - 16
    d.text((PAD, title_y), title, font=f_title, fill=_hex_to_rgb(t['text']))
    f_sub = _load_font(FONT_REG, SUBTITLE_SZ)
    _draw_subtitle_multiline(d, subtitle, PAD, title_y + title_size + 24, f_sub, _hex_to_rgb(t['text2']), max_title_w)
    f_domain = _load_font(FONT_REG, DOM_FONT_SZ)
    db = d.textbbox((0, 0), domain, font=f_domain)
    dw, dh = db[2] - db[0], db[3] - db[1]
    bx_w, bx_h = dw + DOM_PAD_X * 2, dh + DOM_PAD_Y * 2
    bx2, by2 = W - MARGIN, H - MARGIN
    bx1, by1 = bx2 - bx_w, by2 - bx_h
    d.text(((bx1 + bx2) / 2, (by1 + by2) / 2), domain, font=f_domain, fill=_hex_to_rgb(t['accent2']), anchor="mm")
    _draw_bracket(d, bx1, by1, bx2, by2, _hex_to_rgb(t['accent2']))
    return img

@lru_cache(maxsize=300)
def _cached_og(title, subtitle, theme, prompt, domain):
    img = make_og(title, subtitle, theme, prompt, domain)
    buf = io.BytesIO()
    img.save(buf, 'PNG', optimize=True)
    return buf.getvalue()

@og_bp.route('/og-image')
def og_image():
    title, subtitle, theme = request.args.get('title', 'Yiğit Gülyurt')[:80], request.args.get('subtitle', '')[:120], request.args.get('theme', 'default')
    icon, prompt, domain = request.args.get('icon', '')[:20], request.args.get('prompt', '$ whoami')[:60], request.args.get('domain', 'yigitgulyurt.net.tr')[:50]
    try:
        if icon and "\\" in icon: icon = icon.encode('utf-8').decode('unicode_escape').encode('utf-16', 'surrogatepass').decode('utf-16')
    except Exception: pass
    full_prompt = f"{icon} {prompt}".strip() if icon else prompt
    data = _cached_og(title, subtitle, theme, full_prompt, domain)
    resp = send_file(io.BytesIO(data), mimetype='image/png')
    resp.headers['Cache-Control'] = 'public, max-age=3600'
    return resp

# --- Tools Routes ---
@tools_bp.route('/')
def index():
    return render_template('tools/index.html')

@tools_bp.route('/sifre-olusturucu')
def password_gen():
    return render_template('tools/password_gen.html')

@tools_bp.route('/karakter-sayici')
def char_counter():
    return render_template('tools/char_counter.html')

@tools_bp.route('/json-formatlayici')
def json_formatter():
    return render_template('tools/json_formatter.html')

@tools_bp.route('/base64-donusturucu')
def base64_converter():
    return render_template('tools/base64_converter.html')

@tools_bp.route('/ip-bilgisi')
def ip_info():
    user_ip = request.headers.get('X-Forwarded-For', request.remote_addr)
    if ',' in user_ip:
        user_ip = user_ip.split(',')[0].strip()
    return render_template('tools/ip_info.html', user_ip=user_ip)

@tools_bp.route('/ip-log', methods=['POST'])
@limiter.limit("10 per minute")
def log_ip():
    data = request.get_json()
    if not data:
        return jsonify({"error": "No data provided"}), 400
    
    try:
        log = IpLog(
            ip=data.get('ip'),
            city=data.get('city'),
            region=data.get('region'),
            country=data.get('country_name'),
            org=data.get('org'),
            asn=data.get('asn'),
            latitude=data.get('latitude'),
            longitude=data.get('longitude'),
            timezone=data.get('timezone'),
            user_agent=request.user_agent.string,
            full_data=data
        )
        db.session.add(log)
        db.session.commit()
        return jsonify({"status": "success"}), 201
    except Exception as e:
        current_app.logger.error(f"IP Log error: {e}")
        return jsonify({"error": str(e)}), 500

@tools_bp.route('/birim-donusturucu')
def unit_converter():
    return render_template('tools/unit_converter.html')

@tools_bp.route('/identicon-olusturucu')
def identicon():
    return render_template('tools/identicon.html')

@tools_bp.route('/url-encode-decode')
def url_encode_decode():
    return render_template('tools/url_encode_decode.html')

@tools_bp.route('/uuid-uretici')
def uuid_generator():
    return render_template('tools/uuid_generator.html')

@tools_bp.route('/regex-test')
def regex_test():
    return render_template('tools/regex_test.html')

@tools_bp.route('/json-yaml-donusturucu')
def json_yaml_converter():
    return render_template('tools/json_yaml_converter.html')

@tools_bp.route('/jwt-decoder')
def jwt_decoder():
    return render_template('tools/jwt_decoder.html')

@tools_bp.route('/http-header-kontrolcusu')
def http_header_checker():
    return render_template('tools/http_header_checker.html')

@tools_bp.route('/css-minifier-beautifier')
def css_minifier_beautifier():
    return render_template('tools/css_minifier_beautifier.html')

@tools_bp.route('/html-entity-encode-decode')
def html_entity_encode_decode():
    return render_template('tools/html_entity_encode_decode.html')

@tools_bp.route('/tarih-saat-donusturucu')
def datetime_converter():
    return render_template('tools/datetime_converter.html')

@tools_bp.route('/renk-donusturucu')
def color_converter():
    return render_template('tools/color_converter.html')

@tools_bp.route('/hash-uretici')
def hash_generator():
    return render_template('tools/hash_generator.html')

@tools_bp.route('/cron-olusturucu')
def cron_generator():
    return render_template('tools/cron_generator.html')

@tools_bp.route('/markdown-preview')
def markdown_preview():
    return render_template('tools/markdown_preview.html')

@tools_bp.route('/csv-goruntuleyici')
def csv_viewer():
    return render_template('tools/csv_viewer.html')

@tools_bp.route('/dns-sorgulayici')
def dns_lookup():
    return render_template('tools/dns_lookup.html')

@tools_bp.route('/ip-subnet-hesaplayici')
def ip_subnet_calculator():
    return render_template('tools/ip_subnet_calculator.html')

@tools_bp.route('/analog-saat')
def analog_clock():
    return render_template('tools/analog_clock.html')

@tools_bp.route('/digital-saat')
def digital_clock():
    return render_template('tools/digital_clock.html')

@tools_bp.route('/kronometre')
def stopwatch():
    return render_template('tools/stopwatch.html')

@tools_bp.route('/geri-sayim')
def countdown_timer():
    return render_template('tools/countdown_timer.html')

@tools_bp.route('/qr-kod-uretici')
def qr_code_generator():
    return render_template('tools/qr_code_generator.html')

@tools_bp.route('/lorem-ipsum')
def lorem_ipsum():
    return render_template('tools/lorem_ipsum.html')

@tools_bp.route('/rastgele-sayi')
def random_number():
    return render_template('tools/random_number.html')

@tools_bp.route('/sifre-guclugu')
def password_strength():
    return render_template('tools/password_strength.html')

@tools_bp.route('/roma-rakamlari')
def roman_numerals():
    return render_template('tools/roman_numerals.html')

@tools_bp.route('/sicaklik-donusturucu')
def temperature_converter():
    return render_template('tools/temperature_converter.html')

@tools_bp.route('/bahsis-hesaplayici')
def tip_calculator():
    return render_template('tools/tip_calculator.html')

# GitHub-style color palette
GITHUB_COLORS = [
    '#0066FF', '#1192AA', '#008C88', '#11AA77', '#44BB44',
    '#88CC44', '#CCDD44', '#FFDD00', '#FFAA00', '#FF8800',
    '#EE5555', '#DD4477', '#AA3399', '#8833DD', '#6644FF'
]


def hash_string(text):
    """Simple hash function for string"""
    hash_val = 0
    for char in text:
        hash_val = ((hash_val << 5) - hash_val) + ord(char)
        hash_val = hash_val & hash_val
    return abs(hash_val)


def generate_identicon_svg(text, size=420, grid_size=5):
    """Generate GitHub-style identicon SVG"""
    hash_val = hash_string(text)
    color_index = hash_val % len(GITHUB_COLORS)
    fg_color = GITHUB_COLORS[color_index]
    bg_color = '#f0f0f0'
    
    cell_size = size // (grid_size + 1)  # +1 for padding
    padding = (size - cell_size * grid_size) // 2
    
    svg = f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 {size} {size}">'
    svg += f'<rect x="0" y="0" width="{size}" height="{size}" fill="{bg_color}"/>'
    
    for y in range(grid_size):
        for x in range((grid_size + 1) // 2):
            bit_index = y * grid_size + x
            bit_value = (hash_val >> (bit_index % 32)) & 1
            
            if bit_value:
                rect_x = padding + x * cell_size
                rect_y = padding + y * cell_size
                svg += f'<rect x="{rect_x}" y="{rect_y}" width="{cell_size}" height="{cell_size}" fill="{fg_color}"/>'
                
                mirror_x = grid_size - 1 - x
                if mirror_x != x:
                    mirror_rect_x = padding + mirror_x * cell_size
                    svg += f'<rect x="{mirror_rect_x}" y="{rect_y}" width="{cell_size}" height="{cell_size}" fill="{fg_color}"/>'
    
    svg += '</svg>'
    return svg


@tools_bp.route('/api/identicon', methods=['GET', 'POST'])
@limiter.limit("100 per minute")
def identicon_api():
    """API endpoint for generating identicons"""
    if request.method == 'GET':
        text = request.args.get('text', 'default')
        size = int(request.args.get('size', 420))
        grid = int(request.args.get('grid', 5))
    else:  # POST
        data = request.get_json() or {}
        text = data.get('text', 'default')
        size = int(data.get('size', 420))
        grid = int(data.get('grid', 5))
    
    # Validate parameters
    size = max(64, min(1024, size))
    grid = max(3, min(10, grid))
    
    svg = generate_identicon_svg(text, size, grid)
    return Response(svg, mimetype='image/svg+xml')


@tools_bp.route('/api/identicon/png', methods=['GET', 'POST'])
@limiter.limit("50 per minute")
def identicon_png_api():
    """API endpoint for generating identicons as PNG"""
    from io import BytesIO
    from reportlab.graphics import renderPM
    from svglib.svglib import svg2rlg
    
    if request.method == 'GET':
        text = request.args.get('text', 'default')
        size = int(request.args.get('size', 420))
        grid = int(request.args.get('grid', 5))
    else:  # POST
        data = request.get_json() or {}
        text = data.get('text', 'default')
        size = int(data.get('size', 420))
        grid = int(data.get('grid', 5))
    
    # Validate parameters
    size = max(64, min(1024, size))
    grid = max(3, min(10, grid))
    
    svg = generate_identicon_svg(text, size, grid)
    
    try:
        # Convert SVG to PNG using svglib and reportlab
        drawing = svg2rlg(BytesIO(svg.encode('utf-8')))
        png_buffer = BytesIO()
        renderPM.drawToFile(drawing, png_buffer, fmt='PNG', dpi=300)
        png_buffer.seek(0)
        return Response(png_buffer.getvalue(), mimetype='image/png')
    except Exception as e:
        # Fallback: return SVG if conversion fails
        current_app.logger.warning(f"PNG conversion failed: {e}, falling back to SVG")
        return Response(svg, mimetype='image/svg+xml')


@main_bp.route('/dosya-donusturucu', methods=['GET', 'POST'])
@limiter.limit("30 per hour", methods=['POST'])
def file_converter():
    import hashlib
    import time
    import tracemalloc
    import os

    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
    MAX_URL_RESPONSE_SIZE = 50 * 1024 * 1024  # 50 MB
    PBKDF2_ITERATIONS = 100000
    MAX_WORKERS = min(4, os.cpu_count() or 2)  # Paralel işlem için worker sayısı
    CACHE_ENABLED = True
    CACHE_MAX_SIZE = 100  # MB
    CACHE_TTL = 3600  # saniye (1 saat)

    from io import BytesIO
    from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
    from urllib.parse import urlparse
    import zipfile

    # Lazy Loaded libraries (sadece ihtiyaç duyulduğunda yüklenecek)
    _libraries = {}

    def get_library(name):
        if name not in _libraries:
            if name == 'docx':
                from docx import Document
                _libraries['docx'] = Document
            elif name == 'reportlab':
                from reportlab.lib.pagesizes import letter, A4
                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib import colors
                _libraries['reportlab'] = {
                    'A4': A4,
                    'SimpleDocTemplate': SimpleDocTemplate,
                    'Paragraph': Paragraph,
                    'Spacer': Spacer,
                    'Table': Table,
                    'TableStyle': TableStyle,
                    'getSampleStyleSheet': getSampleStyleSheet,
                    'colors': colors
                }
            elif name == 'pypdf2':
                from PyPDF2 import PdfReader, PdfWriter
                _libraries['pypdf2'] = {'PdfReader': PdfReader, 'PdfWriter': PdfWriter}
            elif name == 'openpyxl':
                import openpyxl
                _libraries['openpyxl'] = openpyxl
            elif name == 'pptx':
                from pptx import Presentation
                _libraries['pptx'] = Presentation
            elif name == 'cryptography':
                from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
                from cryptography.hazmat.primitives import hashes
                from cryptography.hazmat.backends import default_backend
                from cryptography.hazmat.primitives.ciphers.aead import AESGCM
                _libraries['cryptography'] = {
                    'PBKDF2HMAC': PBKDF2HMAC,
                    'hashes': hashes,
                    'default_backend': default_backend,
                    'AESGCM': AESGCM
                }
            elif name == 'svglib':
                from svglib.svglib import svg2rlg
                from reportlab.graphics import renderPDF
                _libraries['svglib'] = {'svg2rlg': svg2rlg, 'renderPDF': renderPDF}
            elif name == 'pdf2image':
                from pdf2image import convert_from_bytes
                _libraries['pdf2image'] = convert_from_bytes
            elif name == 'requests':
                import requests
                _libraries['requests'] = requests
        return _libraries[name]

    # Basit önbellek mekanizması
    _cache = {}
    _cache_timestamps = {}
    _cache_sizes = {}

    def get_cache_key(file_content, target_format, operation):
        file_hash = hashlib.sha256(file_content).hexdigest()
        return f"{file_hash}_{target_format}_{operation}"

    def get_from_cache(cache_key):
        if not CACHE_ENABLED:
            return None
        if cache_key in _cache:
            if time.time() - _cache_timestamps[cache_key] < CACHE_TTL:
                current_app.logger.info(f"Cache hit: {cache_key}")
                return _cache[cache_key]
            else:
                del _cache[cache_key]
                del _cache_timestamps[cache_key]
                del _cache_sizes[cache_key]
        return None

    def add_to_cache(cache_key, result):
        if not CACHE_ENABLED:
            return
        result_size = len(result[1]) / (1024 * 1024)
        total_cache_size = sum(_cache_sizes.values())
        if total_cache_size + result_size > CACHE_MAX_SIZE:
            oldest_key = min(_cache_timestamps, key=_cache_timestamps.get)
            del _cache[oldest_key]
            del _cache_timestamps[oldest_key]
            del _cache_sizes[oldest_key]
        _cache[cache_key] = result
        _cache_timestamps[cache_key] = time.time()
        _cache_sizes[cache_key] = result_size
        current_app.logger.info(f"Cache added: {cache_key}, size: {result_size:.2f} MB")

    # Bellek izleme
    def log_memory_usage(step):
        try:
            snapshot = tracemalloc.take_snapshot()
            top_stats = snapshot.statistics('lineno')
            current_app.logger.info(f"[Memory {step}] Top 3 memory usage:")
            for stat in top_stats[:3]:
                current_app.logger.info(f"  {stat}")
        except Exception as e:
            current_app.logger.warning(f"Memory tracking failed: {e}")

        files = []
        target_format = request.form.get('target_format')
        operation = request.form.get('operation', 'convert')

        url_input = request.form.get('url')
        if url_input:
            try:
                parsed = urlparse(url_input)
                if not parsed.scheme or not parsed.netloc:
                    return jsonify({'error': 'Geçersiz URL'}), 400
                
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
                }
                
                with requests.get(url_input, headers=headers, stream=True, timeout=30) as r:
                    r.raise_for_status()
                    
                    content_length = r.headers.get('Content-Length')
                    if content_length and int(content_length) > MAX_URL_RESPONSE_SIZE:
                        return jsonify({'error': 'Dosya çok büyük (max 50 MB)'}), 400
                    
                    content = b''
                    total_read = 0
                    for chunk in r.iter_content(chunk_size=8192):
                        total_read += len(chunk)
                        if total_read > MAX_URL_RESPONSE_SIZE:
                            return jsonify({'error': 'Dosya çok büyük (max 50 MB)'}), 400
                        content += chunk
                    
                    filename = url_input.split('/')[-1]
                    if not filename or '.' not in filename:
                        filename = 'indirilen_dosya'
                        content_type = r.headers.get('Content-Type', '')
                        if 'pdf' in content_type:
                            filename += '.pdf'
                        elif 'image/jpeg' in content_type or 'image/jpg' in content_type:
                            filename += '.jpg'
                        elif 'image/png' in content_type:
                            filename += '.png'
                        elif 'text/plain' in content_type:
                            filename += '.txt'
                        else:
                            filename += '.bin'
                    
                    class FileObj:
                        def __init__(self, filename, content):
                            self.filename = filename
                            self._content = content
                            self._pos = 0
                            self.content_length = len(content)
                        def read(self, size=None):
                            if size is None:
                                result = self._content[self._pos:]
                                self._pos = len(self._content)
                                return result
                            result = self._content[self._pos:self._pos + size]
                            self._pos += len(result)
                            return result
                        def seek(self, pos, whence=0):
                            if whence == 0:
                                self._pos = pos
                            elif whence == 1:
                                self._pos += pos
                            elif whence == 2:
                                self._pos = len(self._content) + pos
                        def tell(self):
                            return self._pos
                    
                    files.append(FileObj(filename, content))
                    
            except Exception as e:
                return jsonify({'error': f'URL\'den dosya indirilemedi: {str(e)}'}), 400
        else:
            files = request.files.getlist('files')

        if not files or len(files) == 0:
            return jsonify({'error': 'Dosya seçilmedi veya URL girilmedi'}), 400

        if operation == 'convert' and not target_format:
            return jsonify({'error': 'Hedef format seçilmedi'}), 400

        try:
            if operation == 'merge_pdf':
                pdf_writer = PdfWriter()
                for file in files:
                    if file.filename == '':
                        continue
                    file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
                    if file_ext != 'pdf':
                        return jsonify({'error': f'{file.filename} PDF değil'}), 400
                    pdf_reader = PdfReader(BytesIO(file.read()))
                    for page in pdf_reader.pages:
                        pdf_writer.add_page(page)
                output_buffer = BytesIO()
                pdf_writer.write(output_buffer)
                output_buffer.seek(0)
                return send_file(
                    output_buffer,
                    as_attachment=True,
                    download_name='birlesmis.pdf',
                    mimetype='application/pdf'
                )
            elif operation == 'encrypt':
                password = request.form.get('password', '')
                if not password or len(password) < 8:
                    return jsonify({'error': 'Şifre en az 8 karakter olmalı'}), 400
                
                def process_encrypt_file(file):
                    if file.filename == '':
                        return None
                    file_content = file.read()
                    salt = os.urandom(16)
                    iv = os.urandom(12)
                    kdf = PBKDF2HMAC(
                        algorithm=hashes.SHA256(),
                        length=32,
                        salt=salt,
                        iterations=PBKDF2_ITERATIONS,
                        backend=default_backend()
                    )
                    key = kdf.derive(password.encode('utf-8'))
                    aesgcm = AESGCM(key)
                    encrypted_data = aesgcm.encrypt(iv, file_content, associated_data=None)
                    final_data = salt + iv + encrypted_data
                    return (file.filename + '.encrypted', final_data)
                
                processed_files = []
                with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                    futures = [executor.submit(process_encrypt_file, file) for file in files]
                    for future in as_completed(futures):
                        result = future.result()
                        if result:
                            processed_files.append(result)
                
                if len(processed_files) == 1:
                    filename, content = processed_files[0]
                    output_buffer = BytesIO(content)
                    output_buffer.seek(0)
                    return send_file(
                        output_buffer,
                        as_attachment=True,
                        download_name=filename,
                        mimetype='application/octet-stream'
                    )
                else:
                    zip_buffer = BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                        for filename, content in processed_files:
                            zipf.writestr(filename, content)
                    zip_buffer.seek(0)
                    return send_file(
                        zip_buffer,
                        as_attachment=True,
                        download_name='sifreli_dosyalar.zip',
                        mimetype='application/zip'
                    )
            elif operation == 'decrypt':
                password = request.form.get('password', '')
                if not password:
                    return jsonify({'error': 'Lütfen şifrenizi girin'}), 400
                
                def process_decrypt_file(file):
                    if file.filename == '':
                        return None
                    file_content = file.read()
                    if len(file_content) < 28:
                        raise Exception(f'{file.filename} geçersiz şifreli dosya')
                    salt = file_content[:16]
                    iv = file_content[16:28]
                    encrypted_data = file_content[28:]
                    kdf = PBKDF2HMAC(
                        algorithm=hashes.SHA256(),
                        length=32,
                        salt=salt,
                        iterations=PBKDF2_ITERATIONS,
                        backend=default_backend()
                    )
                    key = kdf.derive(password.encode('utf-8'))
                    aesgcm = AESGCM(key)
                    decrypted_data = aesgcm.decrypt(iv, encrypted_data, associated_data=None)
                    original_filename = file.filename.replace('.encrypted', '')
                    if original_filename == file.filename:
                        original_filename = f'decrypted_{file.filename}'
                    return (original_filename, decrypted_data)
                
                processed_files = []
                with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                    futures = [executor.submit(process_decrypt_file, file) for file in files]
                    for future in as_completed(futures):
                        try:
                            result = future.result()
                            if result:
                                processed_files.append(result)
                        except Exception:
                            return jsonify({'error': 'Hatalı şifre veya bozuk dosya'}), 400
                
                if len(processed_files) == 1:
                    filename, content = processed_files[0]
                    output_buffer = BytesIO(content)
                    output_buffer.seek(0)
                    return send_file(
                        output_buffer,
                        as_attachment=True,
                        download_name=filename
                    )
                else:
                    zip_buffer = BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                        for filename, content in processed_files:
                            zipf.writestr(filename, content)
                    zip_buffer.seek(0)
                    return send_file(
                        zip_buffer,
                        as_attachment=True,
                        download_name='cozulmus_dosyalar.zip',
                        mimetype='application/zip'
                    )

            def convert_single_file(file):
                if file.filename == '':
                    return None

                file_size = 0
                pre_read_content = None
                if hasattr(file, 'content_length') and file.content_length:
                    file_size = file.content_length
                elif hasattr(file, 'seek') and hasattr(file, 'tell'):
                    file.seek(0, 2)
                    file_size = file.tell()
                    file.seek(0)
                else:
                    pre_read_content = file.read()
                    file_size = len(pre_read_content)
                
                if file_size > MAX_FILE_SIZE:
                    raise Exception(f'{file.filename} çok büyük (max 50 MB)')

                original_filename = secure_filename(file.filename)
                file_ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
                if pre_read_content is not None:
                    file_content = pre_read_content
                else:
                    file_content = file.read()

                is_image = False
                is_svg = False
                is_text = False
                is_docx = False
                is_pdf = False
                is_xlsx = False
                is_pptx = False
                is_video = False
                is_audio = False

                try:
                    img = Image.open(BytesIO(file_content))
                    is_image = True
                except Exception:
                    pass

                if file_ext == 'svg':
                    is_svg = True
                    is_image = True

                text_extensions = ['txt', 'md', 'csv', 'json', 'xml', 'html', 'css', 'js', 'py', 'c', 'cpp', 'java', 'php', 'rb', 'go', 'rs']
                is_text = file_ext in text_extensions

                video_extensions = ['mp4', 'avi', 'mov', 'mkv', 'webm', 'flv', 'wmv', 'm4v', '3gp']
                is_video = file_ext in video_extensions

                audio_extensions = ['mp3', 'wav', 'ogg', 'flac', 'aac', 'm4a', 'wma', 'opus']
                is_audio = file_ext in audio_extensions

                if file_ext == 'docx':
                    is_docx = True
                if file_ext == 'pdf':
                    is_pdf = True
                if file_ext == 'xlsx':
                    is_xlsx = True
                if file_ext == 'pptx':
                    is_pptx = True

                output_buffer = BytesIO()
                output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'

                if target_format == 'pdf':
                    if is_svg:
                        from svglib.svglib import svg2rlg
                        from reportlab.graphics import renderPDF
                        drawing = svg2rlg(BytesIO(file_content))
                        renderPDF.drawToFile(drawing, output_buffer)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    elif is_image:
                        img = Image.open(BytesIO(file_content))
                        if img.mode in ('RGBA', 'P'):
                            img = img.convert('RGB')
                        img.save(output_buffer, format='PDF')
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    elif is_docx:
                        doc = Document(BytesIO(file_content))
                        styles = getSampleStyleSheet()
                        doc_pdf = SimpleDocTemplate(output_buffer, pagesize=A4)
                        story = []
                        for para in doc.paragraphs:
                            story.append(Paragraph(para.text, styles['Normal']))
                            story.append(Spacer(1, 6))
                        for table in doc.tables:
                            data = []
                            for row in table.rows:
                                row_data = []
                                for cell in row.cells:
                                    row_data.append(cell.text)
                                data.append(row_data)
                            if data:
                                t = Table(data)
                                t.setStyle(TableStyle([
                                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                                ]))
                                story.append(t)
                                story.append(Spacer(1, 12))
                        doc_pdf.build(story)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    elif is_xlsx:
                        wb = openpyxl.load_workbook(BytesIO(file_content))
                        styles = getSampleStyleSheet()
                        doc_pdf = SimpleDocTemplate(output_buffer, pagesize=A4)
                        story = []
                        for sheet_name in wb.sheetnames:
                            story.append(Paragraph(sheet_name, styles['Heading1']))
                            story.append(Spacer(1, 12))
                            ws = wb[sheet_name]
                            data = []
                            for row in ws.iter_rows(values_only=True):
                                row_data = [str(cell) if cell is not None else '' for cell in row]
                                data.append(row_data)
                            if data:
                                t = Table(data)
                                t.setStyle(TableStyle([
                                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                                ]))
                                story.append(t)
                                story.append(Spacer(1, 12))
                        doc_pdf.build(story)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    elif is_pptx:
                        prs = Presentation(BytesIO(file_content))
                        styles = getSampleStyleSheet()
                        doc_pdf = SimpleDocTemplate(output_buffer, pagesize=A4)
                        story = []
                        for i, slide in enumerate(prs.slides, 1):
                            story.append(Paragraph(f'Slayt {i}', styles['Heading1']))
                            story.append(Spacer(1, 12))
                            for shape in slide.shapes:
                                if hasattr(shape, 'text') and shape.text.strip():
                                    story.append(Paragraph(shape.text, styles['Normal']))
                                    story.append(Spacer(1, 6))
                            story.append(Spacer(1, 12))
                        doc_pdf.build(story)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    elif is_text or file_ext == 'txt':
                        text = file_content.decode('utf-8', errors='replace')
                        styles = getSampleStyleSheet()
                        doc_pdf = SimpleDocTemplate(output_buffer, pagesize=A4)
                        story = []
                        for line in text.split('\n'):
                            if line.strip():
                                story.append(Paragraph(line, styles['Normal']))
                            else:
                                story.append(Spacer(1, 6))
                        doc_pdf.build(story)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                    else:
                        raise Exception(f'{original_filename} bu dosya türü PDF\'ye dönüştürülemez')

                elif target_format in ['jpg', 'jpeg', 'png', 'webp', 'gif', 'bmp', 'tiff']:
                    if is_svg:
                        from svglib.svglib import svg2rlg
                        from reportlab.graphics import renderPDF
                        from pdf2image import convert_from_bytes
                        import os
                        temp_pdf_buffer = BytesIO()
                        drawing = svg2rlg(BytesIO(file_content))
                        renderPDF.drawToFile(drawing, temp_pdf_buffer)
                        temp_pdf_buffer.seek(0)
                        poppler_path = None
                        if os.path.exists('/usr/bin'):
                            poppler_path = '/usr/bin'
                        elif os.path.exists('/usr/local/bin'):
                            poppler_path = '/usr/local/bin'
                        convert_kwargs = {}
                        if poppler_path:
                            convert_kwargs['poppler_path'] = poppler_path
                        images = convert_from_bytes(temp_pdf_buffer.getvalue(), **convert_kwargs)
                        img = images[0]
                        if target_format in ['jpg', 'jpeg']:
                            if img.mode in ('RGBA', 'P'):
                                img = img.convert('RGB')
                            img.save(output_buffer, format='JPEG', quality=95)
                        else:
                            if target_format == 'webp' and img.mode in ('RGBA', 'P'):
                                img.save(output_buffer, format='WEBP', quality=90)
                            else:
                                img.save(output_buffer, format=target_format.upper())
                        output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'
                    elif is_image:
                        img = Image.open(BytesIO(file_content))
                        if target_format in ['jpg', 'jpeg']:
                            if img.mode in ('RGBA', 'P'):
                                img = img.convert('RGB')
                            img.save(output_buffer, format='JPEG', quality=95)
                        else:
                            if target_format == 'webp' and img.mode in ('RGBA', 'P'):
                                img.save(output_buffer, format='WEBP', quality=90)
                            else:
                                img.save(output_buffer, format=target_format.upper())
                        output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'
                    elif is_pdf:
                        try:
                            from pdf2image import convert_from_bytes
                            import os
                            poppler_path = None
                            if os.path.exists('/usr/bin'):
                                poppler_path = '/usr/bin'
                            elif os.path.exists('/usr/local/bin'):
                                poppler_path = '/usr/local/bin'
                            
                            convert_kwargs = {}
                            if poppler_path:
                                convert_kwargs['poppler_path'] = poppler_path
                            
                            images = convert_from_bytes(file_content, **convert_kwargs)
                            if len(images) == 1:
                                img = images[0]
                                if target_format in ['jpg', 'jpeg']:
                                    img.save(output_buffer, format='JPEG', quality=95)
                                else:
                                    img.save(output_buffer, format=target_format.upper())
                                output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'
                            else:
                                zip_buffer = BytesIO()
                                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                                    for i, img in enumerate(images, 1):
                                        img_buffer = BytesIO()
                                        if target_format in ['jpg', 'jpeg']:
                                            img.save(img_buffer, format='JPEG', quality=95)
                                        else:
                                            img.save(img_buffer, format=target_format.upper())
                                        img_buffer.seek(0)
                                        zipf.writestr(f'{original_filename.rsplit(".", 1)[0]}_sayfa_{i}.{target_format}', img_buffer.getvalue())
                                zip_buffer.seek(0)
                                output_buffer = zip_buffer
                                output_filename = original_filename.rsplit('.', 1)[0] + '_sayfalar.zip'
                        except Exception as e:
                            import os
                            poppler_paths = ['/usr/bin', '/usr/local/bin']
                            available_paths = [p for p in poppler_paths if os.path.exists(p)]
                            path_info = f" (Poppler yolları kontrol edildi: {', '.join(available_paths) if available_paths else 'yol bulunamadı'})"
                            raise Exception(f'PDF\'den resim dönüşümü için poppler yüklü ve PATH\'de olmalı. Hata: {str(e)}{path_info}')
                    else:
                        raise Exception(f'{original_filename} sadece resim ve PDF dosyaları resim formatlarına dönüştürülebilir')

                elif target_format == 'txt':
                    if is_docx:
                        doc = Document(BytesIO(file_content))
                        text_content = '\n'.join([para.text for para in doc.paragraphs])
                        output_buffer.write(text_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.txt'
                    elif is_pdf:
                        pdf_reader = PdfReader(BytesIO(file_content))
                        text_content = ''
                        for page in pdf_reader.pages:
                            text_content += page.extract_text() + '\n'
                        output_buffer.write(text_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.txt'
                    elif is_xlsx:
                        wb = openpyxl.load_workbook(BytesIO(file_content))
                        text_content = ''
                        for sheet_name in wb.sheetnames:
                            text_content += f'--- {sheet_name} ---\n'
                            ws = wb[sheet_name]
                            for row in ws.iter_rows(values_only=True):
                                row_data = [str(cell) if cell is not None else '' for cell in row]
                                text_content += '\t'.join(row_data) + '\n'
                            text_content += '\n'
                        output_buffer.write(text_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.txt'
                    elif is_text or file_ext == 'txt':
                        output_buffer.write(file_content)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.txt'
                    else:
                        raise Exception(f'{original_filename} bu dosya türü TXT\'ye dönüştürülemez')

                elif target_format == 'csv':
                    if is_xlsx:
                        wb = openpyxl.load_workbook(BytesIO(file_content))
                        if len(wb.sheetnames) == 1:
                            ws = wb.active
                            text_content = ''
                            for row in ws.iter_rows(values_only=True):
                                row_data = [str(cell) if cell is not None else '' for cell in row]
                                text_content += ','.join(row_data) + '\n'
                            output_buffer.write(text_content.encode('utf-8'))
                            output_filename = original_filename.rsplit('.', 1)[0] + '.csv'
                        else:
                            zip_buffer = BytesIO()
                            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                                for sheet_name in wb.sheetnames:
                                    ws = wb[sheet_name]
                                    text_content = ''
                                    for row in ws.iter_rows(values_only=True):
                                        row_data = [str(cell) if cell is not None else '' for cell in row]
                                        text_content += ','.join(row_data) + '\n'
                                    zipf.writestr(f'{original_filename.rsplit(".", 1)[0]}_{sheet_name}.csv', text_content.encode('utf-8'))
                            zip_buffer.seek(0)
                            output_buffer = zip_buffer
                            output_filename = original_filename.rsplit('.', 1)[0] + '_sayfalar.zip'
                    else:
                        raise Exception(f'{original_filename} sadece Excel dosyaları CSV\'ye dönüştürülebilir')

                elif target_format == 'html':
                    if is_docx:
                        doc = Document(BytesIO(file_content))
                        html_content = '<!DOCTYPE html><html><head><meta charset="UTF-8"><title></title></head><body>'
                        for para in doc.paragraphs:
                            html_content += f'<p>{para.text}</p>'
                        html_content += '</body></html>'
                        output_buffer.write(html_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.html'
                    elif is_text or file_ext == 'txt':
                        text = file_content.decode('utf-8', errors='replace')
                        html_content = '<!DOCTYPE html><html><head><meta charset="UTF-8"><title></title></head><body>'
                        for line in text.split('\n'):
                            if line.strip():
                                html_content += f'<p>{line}</p>'
                            else:
                                html_content += '<br>'
                        html_content += '</body></html>'
                        output_buffer.write(html_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.html'
                    else:
                        raise Exception(f'{original_filename} bu dosya türü HTML\'ye dönüştürülemez')

                elif target_format == 'md':
                    if is_docx:
                        doc = Document(BytesIO(file_content))
                        md_content = ''
                        for para in doc.paragraphs:
                            md_content += para.text + '\n\n'
                        output_buffer.write(md_content.encode('utf-8'))
                        output_filename = original_filename.rsplit('.', 1)[0] + '.md'
                    elif is_text or file_ext == 'txt':
                        output_buffer.write(file_content)
                        output_filename = original_filename.rsplit('.', 1)[0] + '.md'
                    else:
                        raise Exception(f'{original_filename} bu dosya türü Markdown\'a dönüştürülemez')

                elif target_format in ['mp4', 'avi', 'mov', 'mkv', 'webm', 'flv', 'wmv', 'm4v', '3gp']:
                    if is_video:
                        import tempfile
                        import os
                        import subprocess
                        try:
                            with tempfile.NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_in:
                                temp_in.write(file_content)
                                temp_in_path = temp_in.name
                            
                            temp_out_path = tempfile.mktemp(suffix=f'.{target_format}')
                            
                            ffmpeg_path = '/usr/bin/ffmpeg'
                            if not os.path.exists(ffmpeg_path):
                                ffmpeg_path = 'ffmpeg'
                            
                            cmd = [ffmpeg_path, '-i', temp_in_path, '-y', temp_out_path]
                            result = subprocess.run(cmd, capture_output=True, text=True)
                            
                            if result.returncode != 0:
                                raise Exception(f'FFmpeg hatası: {result.stderr}')
                            
                            with open(temp_out_path, 'rb') as f:
                                output_buffer.write(f.read())
                            
                            output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'
                            
                            os.unlink(temp_in_path)
                            os.unlink(temp_out_path)
                        except Exception as e:
                            raise Exception(f'Video dönüştürme hatası: {str(e)}. FFmpeg yüklü mü?')
                    else:
                        raise Exception(f'{original_filename} sadece video dosyaları video formatlarına dönüştürülebilir')

                elif target_format in ['mp3', 'wav', 'ogg', 'flac', 'aac', 'm4a', 'wma', 'opus']:
                    if is_audio or is_video:
                        import tempfile
                        import os
                        import subprocess
                        try:
                            input_ext = file_ext
                            with tempfile.NamedTemporaryFile(suffix=f'.{input_ext}', delete=False) as temp_in:
                                temp_in.write(file_content)
                                temp_in_path = temp_in.name
                            
                            temp_out_path = tempfile.mktemp(suffix=f'.{target_format}')
                            
                            ffmpeg_path = '/usr/bin/ffmpeg'
                            if not os.path.exists(ffmpeg_path):
                                ffmpeg_path = 'ffmpeg'
                            
                            cmd = [ffmpeg_path, '-i', temp_in_path, '-y', temp_out_path]
                            result = subprocess.run(cmd, capture_output=True, text=True)
                            
                            if result.returncode != 0:
                                raise Exception(f'FFmpeg hatası: {result.stderr}')
                            
                            with open(temp_out_path, 'rb') as f:
                                output_buffer.write(f.read())
                            
                            output_filename = original_filename.rsplit('.', 1)[0] + f'.{target_format}'
                            
                            os.unlink(temp_in_path)
                            os.unlink(temp_out_path)
                        except Exception as e:
                            raise Exception(f'Ses dönüştürme hatası: {str(e)}. FFmpeg yüklü mü?')
                    else:
                        raise Exception(f'{original_filename} sadece ses ve video dosyaları ses formatlarına dönüştürülebilir')

                else:
                    raise Exception('Desteklenmeyen hedef format')

                output_buffer.seek(0)
                return (output_filename, output_buffer.getvalue())

            converted_files = []
            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                futures = [executor.submit(convert_single_file, file) for file in files]
                for future in as_completed(futures):
                    try:
                        result = future.result()
                        if result:
                            converted_files.append(result)
                    except Exception as e:
                        current_app.logger.error(f"Dosya dönüştürme hatası: {e}")
                        return jsonify({'error': f'Dönüştürme sırasında bir hata oluştu: {str(e)}'}), 500

        except Exception as e:
            current_app.logger.error(f"Dosya dönüştürme hatası: {e}")
            return jsonify({'error': f'Dönüştürme sırasında bir hata oluştu: {str(e)}'}), 500

        if len(converted_files) == 1:
            filename, content = converted_files[0]
            return send_file(
                BytesIO(content),
                as_attachment=True,
                download_name=filename,
                mimetype='application/octet-stream'
            )
        else:
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for filename, content in converted_files:
                    zipf.writestr(filename, content)
            zip_buffer.seek(0)
            return send_file(
                zip_buffer,
                as_attachment=True,
                download_name='donusturulmus_dosyalar.zip',
                mimetype='application/zip'
            )

    return render_template('tools/file_converter.html')
